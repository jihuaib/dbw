# -*- coding: utf-8 -*-
"""生成答辩 PPT：docs/DetOps-答辩.pptx（python-pptx，16:9）。

运行：.venv/bin/python docs/build_deck.py
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x14, 0x2B, 0x4D)
INK = RGBColor(0x1F, 0x23, 0x29)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x0F, 0x76, 0xC7)
ACCENT2 = RGBColor(0xD9, 0x7B, 0x1B)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PALE = RGBColor(0xF3, 0xF6, 0xFA)
LINE = RGBColor(0xD6, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "PingFang SC"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _tf(shape, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=None):
    tf = shape.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(size), bold, color, FONT
    return tf


def box(slide, x, y, w, h, text="", fill=PALE, line=LINE, size=13, bold=False, color=INK,
        align=PP_ALIGN.CENTER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, anchor=MSO_ANCHOR.MIDDLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if text:
        _tf(s, text, size=size, bold=bold, color=color, align=align, anchor=anchor)
        s.text_frame.margin_left = s.text_frame.margin_right = Inches(0.08)
    return s


def arrow(slide, x1, y1, x2, y2, color=MUTED):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(1.5)
    ln = c.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return c


def text(slide, x, y, w, h, lines, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = _tf(tb, lines, size=size, bold=bold, color=color, align=align)
    for p in tf.paragraphs:
        p.space_after = Pt(spacing)
    return tb


def bullets(slide, x, y, w, h, items, size=16):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level = 0
        if isinstance(item, tuple):
            item, level = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = ("▪ " if level == 0 else "– ") + item
        r.font.size = Pt(size - 2 * level)
        r.font.color.rgb = INK if level == 0 else MUTED
        r.font.name = FONT
    return tb


def header(slide, title, sub=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    text(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.6), title, size=26, color=WHITE, bold=True)
    if sub:
        text(slide, Inches(0.62), Inches(0.62), Inches(11.5), Inches(0.4), sub, size=13, color=RGBColor(0xC9, 0xD6, 0xE8))
    foot = text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
                "DetOps · H3C 赛题 6.4.3 AI 大模型在运维场景的诊断一致性", size=10, color=MUTED)
    return slide


def slide(title, sub=""):
    s = prs.slides.add_slide(BLANK)
    return header(s, title, sub)


def num(slide_):
    n = len(prs.slides)
    text(slide_, Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.3), str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ── 1 封面 ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5), "H3C 赛题 6.4.3", size=18, color=RGBColor(0xC9, 0xD6, 0xE8))
text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.4), "AI 大模型在运维场景的诊断一致性", size=40, color=WHITE, bold=True)
text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.8),
     "DetOps —— 一套基于 Agent 运行的运维一致性兜底策略", size=22, color=RGBColor(0xE8, 0xEE, 0xF6))
text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.4), [
    "核心命题：大模型对同一输入无法确定性输出 → 让「输入」确定，让「答案」只生成一次并冻结",
    "交付：真机验证的 60% / 80% / 100% 三级判据全部达标（SSR 100%，字节一致）",
], size=16, color=RGBColor(0xC9, 0xD6, 0xE8))
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.4), "答辩汇报 · 2026-08", size=13, color=RGBColor(0x9F, 0xB3, 0xCC))

# ── 2 题目理解 ────────────────────────────────────────────────────────
s = slide("题目理解：三条设计要求，三级完成度判据", "设备侧不能大改，Agent 要有一致性兜底，多设备要能联动")
box(s, Inches(0.6), Inches(1.4), Inches(5.9), Inches(0.5), "设计要求", fill=NAVY, color=WHITE, bold=True, size=14)
bullets(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(3.2), [
    "① 一套基于 Agent 运行的运维一致性兜底策略",
    "② 设备侧既有能力少变动或不变动",
    "③ 多台设备互联场景的联动分析能力",
], size=17)
box(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.5), "完成度判据（单会话多次 + 多会话，输出一致）", fill=NAVY, color=WHITE, bold=True, size=14)
for i, (pct, desc) in enumerate([("60%", "单一设备、确定的网络故障"),
                                  ("80%", "单一设备、多异常表项故障"),
                                  ("100%", "≥3 台 SPINE-LEAF 三层组网、多异常表项故障")]):
    y = Inches(2.05 + i * 0.95)
    box(s, Inches(6.9), y, Inches(1.2), Inches(0.75), pct, fill=ACCENT, color=WHITE, bold=True, size=20, line=None)
    box(s, Inches(8.2), y, Inches(4.6), Inches(0.75), desc, fill=PALE, size=14, align=PP_ALIGN.LEFT)
box(s, Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.2),
    "关键难点：不是「让模型稳定」（做不到），而是把模型的随机性从「重复提问」这条路径上移出去。\n"
    "我们的回答：同一问题 + 同一设备状态 → 只真正调用一次模型 → 答案按「诊断指纹」冻结，之后逐字节复用。",
    fill=RGBColor(0xFF, 0xF6, 0xEB), line=ACCENT2, size=15, align=PP_ALIGN.LEFT)
num(s)

# ── 3 核心洞察 ────────────────────────────────────────────────────────
s = slide("核心洞察：一致性来自「输入确定 + 答案冻结」，不来自「模型稳定」")
box(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(0.9),
    "诊断指纹 = SHA256( 归一化提问 ‖ 证据快照哈希 ‖ 模型身份 ‖ 命令清单 ‖ 提示词版本 ‖ 归一化版本 ‖ 会话前缀 ‖ 事件摘要 )",
    fill=NAVY, color=WHITE, size=15, bold=True)
cols = [
    ("让输入确定", ["命令只能从手册闭集里选，参数不许猜", "回显归一化：规则擦除 + 表头列擦除 + 实测标定",
                   "事件（syslog/trap）只记类型存在性", "观测行为自身的回声全部排除"]),
    ("只调一次模型", ["Agent 每轮模型调用按「语义键」缓存", "同一段对话第二次出现直接回放",
                     "指纹命中 → F0：零模型调用，字节一致", "设备状态变了 → 指纹变 → 理应重诊"]),
    ("兜不住也有底", ["F2 结构校验、F3 自洽投票", "F4 模型不可用 → 只做证据陈述",
                     "F5 采集缺失 → 列出缺口、下调置信度", "任何版本锚变化 → 旧冻结自动失效"]),
]
for i, (t, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.15)
    box(s, x, Inches(2.6), Inches(3.9), Inches(0.5), t, fill=ACCENT, color=WHITE, bold=True, size=15, line=None)
    bullets(s, x, Inches(3.2), Inches(3.9), Inches(3.4), items, size=14)
num(s)

# ── 4 总体架构 ────────────────────────────────────────────────────────
s = slide("总体架构", "Vue 3 + NetNexusUI 前端 · FastAPI 后端 · LangGraph Agent · 设备只走既有 CLI / syslog / trap")
# 前端
box(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(0.75),
    "前端（Vue 3 / NetNexusUI）：诊断对话 · 知识库 · 设备与拓扑（LLDP 拓扑、网页终端）· Syslog · SNMP（MIB/Trap）· 一致性验证 · 设置弹窗",
    fill=PALE, size=13)
arrow(s, Inches(6.7), Inches(2.15), Inches(6.7), Inches(2.55))
# 后端模块
mods = [("diagnose", "Agent 循环 / 指纹冻结\n六级兜底 / 一致性验证"),
        ("kb", "手册导入（多格式）\n只读命令闭集"),
        ("devices", "SSH/Telnet 接入\n能力探测 / 拓扑 / 终端"),
        ("collect", "采集纪元 / 归一化\n实测标定 / 证据快照"),
        ("events", "Syslog 服务器\npysnmp Trap 接收"),
        ("mibs", "pysmi 编译\nOID 索引 / 树 / 解码")]
for i, (name, desc) in enumerate(mods):
    x = Inches(0.6 + i * 2.05)
    box(s, x, Inches(2.6), Inches(1.9), Inches(0.42), name, fill=ACCENT, color=WHITE, bold=True, size=13, line=None)
    box(s, x, Inches(3.02), Inches(1.9), Inches(0.95), desc, fill=WHITE, size=11)
box(s, Inches(0.6), Inches(4.1), Inches(12.2), Inches(0.45),
    "core：LLM 多服务商适配（Anthropic / DeepSeek / GLM / Qwen / 任意 OpenAI 兼容）· 内容哈希缓存 · SQLite · 版本锚",
    fill=RGBColor(0xE8, 0xF0, 0xF9), size=12)
arrow(s, Inches(3.0), Inches(4.55), Inches(3.0), Inches(5.0))
arrow(s, Inches(10.3), Inches(4.55), Inches(10.3), Inches(5.0))
box(s, Inches(0.6), Inches(5.05), Inches(5.6), Inches(1.5),
    "网络设备（真机 / CNetNexus 实验环境）\n只用既有能力：只读 CLI（SSH/Telnet）、syslog server、snmp trap server\n设备侧零改造 —— 设计要求 ②",
    fill=RGBColor(0xEE, 0xF7, 0xF0), line=GREEN, size=13)
box(s, Inches(7.2), Inches(5.05), Inches(5.6), Inches(1.5),
    "大模型（可换）\n温度 0 + 语义键缓存 + 结构化/固定结构输出\n模型身份进指纹：换模型 = 换口径 = 重诊",
    fill=RGBColor(0xFF, 0xF6, 0xEB), line=ACCENT2, size=13)
num(s)

# ── 5 诊断流程 ────────────────────────────────────────────────────────
s = slide("一次诊断的完整流程", "从提问到冻结答案，每一步都为「同一输入」服务")
steps = [("① 提取信息", "归一化提问\n抽取 IP/接口/设备\n带入会话前缀与事件摘要"),
         ("② Agent 取证", "LangGraph 多轮循环\nrun_cli 工具：闭集校验\n参数不许猜、并行调用"),
         ("③ 采集", "SSH/Telnet 逐台下发\n多轮追加进同一纪元\n(设备, 命令) 全序"),
         ("④ 归一化", "规则擦除 / 表头列擦除\n实测标定兜底\n→ 证据快照哈希"),
         ("⑤ 指纹", "8 个分量 SHA256\n任一分量变即失效"),
         ("⑥ 冻结/复用", "命中 → 原样返回（F0）\n未命中 → 结论 → 冻结")]
for i, (t, d) in enumerate(steps):
    x = Inches(0.6 + i * 2.08)
    box(s, x, Inches(1.5), Inches(1.9), Inches(0.5), t, fill=ACCENT if i != 5 else GREEN, color=WHITE, bold=True, size=13, line=None)
    box(s, x, Inches(2.0), Inches(1.9), Inches(1.5), d, fill=WHITE, size=11)
    if i < 5:
        arrow(s, x + Inches(1.9), Inches(2.75), x + Inches(2.08), Inches(2.75))
box(s, Inches(0.6), Inches(3.8), Inches(12.2), Inches(0.45),
    "全程可观测：对话流实时显示每条工具调用（设备 # 命令）与完整回显、逐 token 流出的结论；结束后可查看送给模型的完整上下文与原始回复",
    fill=PALE, size=12)
bullets(s, Inches(0.6), Inches(4.4), Inches(12.2), Inches(2.5), [
    "指纹命中也要真实重新采集：F0 省的是模型调用，不是采集 —— 状态变了必须能发现（接口 up 之后不会再答「down」）",
    "会话历史进指纹：追问带上下文且确定；重复提问不带自己的历史，单会话连问 N 次仍是同一指纹",
    "被归一化擦掉的时变量（计数器、倒计时）变化不触发重诊 —— 它们本来就不是诊断依据",
], size=14)
num(s)

# ── 6 输入确定性：知识库 ────────────────────────────────────────────
s = slide("让输入确定 ①：命令闭集 —— 手册导入 × 实机探测", "模型只能从手册里选命令；手册说有、这台设备未必有")
bullets(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.3), [
    "手册导入：Word / Markdown，支持文件、文件夹递归、服务器目录递归",
    ("多种写法：表格、标题即命令、【命令】标签、厂商语法行 [ ] { | }、代码块、示例行", 1),
    ("大批量：逐个落盘入库、进度轮询、同内容去重，结果不驻内存", 1),
    ("规则识别不了的格式 → 自动回退 AI 提取", 1),
    "只留只读命令（show/display），危险词过滤；CLI 会话/审计类命令排除（可配置）",
    "必需参数识别：<param> 与厂商命名习惯（vlan-id / interface-number）→ 不下发光杆命令",
    "能力探测：把清单在每台设备上跑一遍，记录「设备 → 不支持」，编排时不再提议",
    "命令清单指纹进诊断指纹：清单变了（比如补进 show current-configuration）→ 重诊",
], size=14)
box(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5), "真机上踩到的坑（都变成了机制）", fill=NAVY, color=WHITE, bold=True, size=14)
bullets(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.6), [
    "命令被小写化 → GE-1 变 ge-1 被拒 → 保留大小写、匹配用小写键",
    "必需参数被当可选剥掉 → Incomplete command 污染证据 → 必需参数不放行",
    "模型凭空猜参数（tag=default）→ 提示词硬规则 + 校验",
    "配置查看命令没导进来 → 接口 shutdown 只能诊断成「物理断链」→ 导入器补格式后根因直指 shutdown",
    "show cli history 这类「观测者的镜子」进了证据 → 永不一致 → 排除",
], size=13)
num(s)

# ── 7 归一化 ──────────────────────────────────────────────────────────
s = slide("让输入确定 ②：归一化 —— 用测量代替猜测", "同一设备状态的两次回显必须归一成同一串字节")
cols = [("规则表（快路径）", ["时间戳 / uptime / 倒计时", "ARP/MAC 老化、OSPF Dead、BGP 报文计数", "LLDP TTL「19 sec remaining」", "光功率等连续量 → 分档离散化"]),
        ("按表头识别易变列", ["找表头 → 按整列空白定列边界", "表头命中 TTL/Age/Dead/Uptime… 整列擦除", "换厂商、换命令不用补规则", "列宽归一：尾随空格漂移不改哈希"]),
        ("实测标定（兜底）", ["同一命令多采几次 → token 级比对", "真正在变的位置冻结成 profile", "profile 并集合并：见过在变就永远擦", "格式无关、厂商无关；换设备重跑一次"])]
for i, (t, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.15)
    box(s, x, Inches(1.4), Inches(3.9), Inches(0.5), t, fill=ACCENT, color=WHITE, bold=True, size=15, line=None)
    bullets(s, x, Inches(2.0), Inches(3.9), Inches(3.0), items, size=13)
box(s, Inches(0.6), Inches(5.2), Inches(12.2), Inches(1.4),
    "为什么必须「实测」：逐条写正则去猜哪一列会变，永远追不上——80% 判据环境首跑 SSR 0%，\n"
    "漏掉的是 LLDP 详情里的倒计时和此前恰好没动过的计数器；重新标定后 SSR 100%。\n"
    "证明链：轮次之间原始回显真实漂移 6～7 处，归一化后快照逐字节相同 —— 一致性是兜出来的，不是输入没动。",
    fill=RGBColor(0xFF, 0xF6, 0xEB), line=ACCENT2, size=13, align=PP_ALIGN.LEFT)
num(s)

# ── 8 Agent ───────────────────────────────────────────────────────────
s = slide("Agent：开源 LangGraph 跑循环，笼子是我们的", "主流 Agent 的体验 + 可重放的确定性")
box(s, Inches(0.6), Inches(1.4), Inches(5.9), Inches(0.5), "框架负责", fill=ACCENT, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(2.6), [
    "ReAct 循环：看一轮回显再决定继续取证还是下结论",
    "同一轮并行多条工具调用，逐 token 流式输出",
    "递归上限保证终止",
], size=14)
box(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.5), "run_cli 工具外的笼子", fill=ACCENT2, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(6.9), Inches(2.0), Inches(5.9), Inches(2.6), [
    "命令必须来自清单、必需参数必须补齐、探明不支持的不放行",
    "跨轮去重、总量预算、设备白名单",
    "先排除「配置使然」再下物理结论（通用取证原则）",
], size=14)
box(s, Inches(0.6), Inches(4.7), Inches(12.2), Inches(0.5), "确定性：语义键缓存", fill=NAVY, color=WHITE, bold=True, size=14)
bullets(s, Inches(0.6), Inches(5.3), Inches(12.2), Inches(1.6), [
    "temperature=0 + SQLite 精确缓存：同一段对话（含工具回显）第二次出现直接回放，模型随机性移出重放路径",
    "踩坑：框架给每条消息随机 uuid、AIMessage 带 run_id/token 用量，混进缓存键后永远命不中（SSR 一度 0%）→ 进键前剥掉，只留角色/内容/工具调用",
    "结果：同一问题 + 同一设备状态 → 同一工具调用序列 → 同一快照 → 同一指纹",
], size=13)
num(s)

# ── 9 兜底阶梯 ────────────────────────────────────────────────────────
s = slide("六级一致性兜底阶梯（设计要求 ① 的交付物）")
rungs = [("F0", "指纹冻结", "指纹命中", "原样返回，零模型调用，字节一致", GREEN),
         ("F1", "快照复用", "同一快照不同问法", "复用同一份事实", ACCENT),
         ("F2", "结构校验", "输出不合结构", "固定次数重试，不放行脏输出", ACCENT),
         ("F3", "自洽投票", "首次生成", "k 次采样按根因取多数票（首答质量保险）", ACCENT),
         ("F4", "模型兜底", "无 Key / 超时 / 报错", "降级为「证据陈述」：只列采到什么，不下结论", ACCENT2),
         ("F5", "缺证兜底", "采集失败", "缺失显式进快照，列出缺口并下调置信度", ACCENT2)]
for i, (lv, name, trig, guar, col) in enumerate(rungs):
    y = Inches(1.4 + i * 0.85)
    box(s, Inches(0.6), y, Inches(0.9), Inches(0.7), lv, fill=col, color=WHITE, bold=True, size=16, line=None)
    box(s, Inches(1.6), y, Inches(2.0), Inches(0.7), name, fill=PALE, bold=True, size=14)
    box(s, Inches(3.7), y, Inches(3.0), Inches(0.7), trig, fill=WHITE, size=13)
    box(s, Inches(6.8), y, Inches(6.0), Inches(0.7), guar, fill=WHITE, size=13, align=PP_ALIGN.LEFT)
text(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.5),
     "越靠上越「硬」：F0 是判据要求的字节一致；F4/F5 保证系统在最坏情况下也给出确定、可解释的输出而不是随机猜。", size=13, color=MUTED)
num(s)

# ── 10 多设备联动 ────────────────────────────────────────────────────
s = slide("多设备联动分析（设计要求 ③）", "拓扑作上下文，跨设备取证，健康设备当对照")
bullets(s, Inches(0.6), Inches(1.4), Inches(6.2), Inches(5.2), [
    "一键 LLDP 拓扑发现：AI 解析 + 正则兜底，自环过滤，双向确认；不依赖厂商",
    "拓扑上下文随提问送给 Agent，并进证据快照",
    "跨设备取证：LEAF2 邻居丢了 → 主动去 SPINE1 看对端口、拿健康的 LEAF1 做对照",
    "100% 场景实测：三台设备各埋一处不同层级的异常（接口配置关闭 / 黑洞静态路由 / OSPF 定时器不匹配），三处全部找齐并分清层级与处置",
    "拓扑页：缩放平移、拖动节点、右键设备 → 网页终端（WebSocket 桥到 telnet/SSH）/ 测试 / 编辑",
], size=14)
box(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.5), "设备主动上报也进诊断", fill=NAVY, color=WHITE, bold=True, size=14)
bullets(s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(4.6), [
    "Syslog 服务器 + pysnmp Trap 接收（只收，不做任何 SNMP 下发）",
    "MIB：源文件导入（自带标准 MIB + 厂商 MIB）→ pysmi 编译 → OID 索引 → 树形浏览 → trap 符号化解码",
    "事件进指纹的三条规矩：只记类型存在性不记次数；CLI 审计事件排除；观测者的镜子命令不进闭集",
    "「刚才发生了什么」这类只有事件能答的问题：模型把 link-up → 邻居建立 → down 的时间线与当前配置对上",
], size=13)
num(s)

# ── 11 实测判据 ──────────────────────────────────────────────────────
s = slide("完成度判据：真机实测全部达标", "1×SPINE + 3×LEAF 真机容器，OSPF 三层组网；每级跑「单会话多次 + 多会话」各 3 轮，轮间真实重新采集")
hdr = ["判据", "故障环境", "快照数", "指纹数", "字节一致", "SSR", "原始回显漂移"]
rows = [["60%", "LEAF2 GE-1 shutdown", "1", "1", "是", "100%", "有"],
        ["80%", "LEAF2：hello 不匹配 + 2 条坏静态路由表项", "1", "1", "是", "100%", "6～7 处"],
        ["100%", "LEAF2 shutdown + SPINE1 黑洞静态路由 + LEAF3 hello 不匹配", "1", "1", "是", "100%", "6 处"]]
widths = [1.0, 4.6, 1.0, 1.0, 1.1, 1.0, 1.5]
x0, y0 = Inches(0.6), Inches(1.45)
x = x0
for wdt, h_ in zip(widths, hdr):
    box(s, x, y0, Inches(wdt), Inches(0.5), h_, fill=NAVY, color=WHITE, bold=True, size=13, shape=MSO_SHAPE.RECTANGLE)
    x += Inches(wdt)
for r_i, row in enumerate(rows):
    x = x0
    for c_i, (wdt, val) in enumerate(zip(widths, row)):
        fill = WHITE if r_i % 2 == 0 else PALE
        col = GREEN if val in ("是", "100%") else INK
        box(s, x, y0 + Inches(0.5 + r_i * 0.62), Inches(wdt), Inches(0.62), val, fill=fill, size=13,
            bold=(c_i == 0 or col == GREEN), color=col, shape=MSO_SHAPE.RECTANGLE)
        x += Inches(wdt)
bullets(s, Inches(0.6), Inches(4.0), Inches(12.2), Inches(2.8), [
    "「原始回显漂移 N 处」是证明链的关键：轮次之间设备回显确实变了，归一化后快照仍逐字节相同",
    "故障用真实配置错误在真实协议栈里种出真实的异常表项（错误静态路由进 RIB/FIB、OSPF 因参数校验拒绝邻居、LSA 自然老化），不是预置的假回显",
    "状态变化实测：shutdown 时答「配置关闭」，no shutdown 后同一会话同一问题指纹改变、重新诊断 —— 不会守着旧答案",
    "所有故障注入脚本化（scripts/lab.sh fault / fault-multi / fault-fabric），可复现",
], size=13)
num(s)

# ── 12 踩坑 ──────────────────────────────────────────────────────────
s = slide("真机逼出来的问题（模拟器一个都发现不了）", "每一条都变成了机制或测试")
left = ["读到提示符前就停 → 空回显 → 判缺失而非有效数据", "命令/回显错位 → 发送前清缓冲",
        "LLDP 自环、列宽漂移、LSA Age → 表头列擦除 + 标定", "首答随机 uuid 进缓存键 → 语义键",
        "事件计数分桶中途翻转 → 只记存在性", "采集自身产生的 CLI 审计日志挤掉状态事件 → SQL 层排除"]
right = ["Docker Desktop NAT 把源 IP 揉成 127.0.0.1 → 每设备独立端口即身份",
         "Linux 上 host.docker.internal 的 UDP 不通 → 实测选址（探测包谁能到用谁）",
         "Python 3.12 移除 asyncore/imp → 迁到 pysnmp 7", "IF-MIB 与 RFC1213 同 OID → SMIv2 优先",
         "H3C 手册格式全漏 → 导入器覆盖标题/标签/语法行/代码块", "厂商预设/排除命令写死 → vendors.json / 设置项"]
box(s, Inches(0.6), Inches(1.4), Inches(5.9), Inches(0.5), "采集与一致性", fill=ACCENT, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.6), left, size=13)
box(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.5), "部署与通用性", fill=ACCENT2, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(6.9), Inches(2.0), Inches(5.9), Inches(4.6), right, size=13)
num(s)

# ── 13 工程化 ────────────────────────────────────────────────────────
s = slide("工程化：测试体系、部署、通用性", "与 CI 同一套用例，不写一次性验证脚本")
cols = [("测试", ["65 个单元测试（离线）：归一化幂等与标定、指纹分量隔离、闭集校验、回显判定、Agent 笼子（假模型）、缓存键、事件、MIB、telnet 桥",
                 "真机集成（-m live）：Agent 问答、F0 冻结、观测接口", "19 步 Playwright 端到端：跑在真机上", "GitHub Actions：单元 + MIB 编译 + 前端构建"]),
        ("部署", ["scripts/start.sh prod：构建前端，后台 0.0.0.0:8099 单端口托管页面/API/WebSocket", "scripts/lab.sh：实验环境起停、故障注入、设备注册、上报链路自检",
                  "Python 3.9～3.13、Linux/macOS；Windows 浏览器访问", "敏感信息不入库：Key/凭据在 data/，备份在仓库外"]),
        ("通用性（无写死）", ["厂商预设 vendors.json：接入协议、关分屏/LLDP 命令、上报命令模板", "命令清单来自导入的手册，MIB 来自导入的源文件",
                            "证据排除命令是设置项", "真机按源 IP 归属事件；实验环境逻辑只在存在 nn-mgmt 网络时启用"])]
for i, (t, items) in enumerate(cols):
    x = Inches(0.6 + i * 4.15)
    box(s, x, Inches(1.4), Inches(3.9), Inches(0.5), t, fill=NAVY, color=WHITE, bold=True, size=15)
    bullets(s, x, Inches(2.0), Inches(3.9), Inches(4.8), items, size=12)
num(s)

# ── 14 总结 ──────────────────────────────────────────────────────────
s = slide("总结与边界")
box(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(1.3),
    "一句话：手册定义可选范围（闭集）→ 实测收窄范围（能力探测）→ 模型在范围内按场景逐轮取证 →\n"
    "校验层保证每条命令合法 → 归一化保证同一状态同一字节 → 指纹冻结保证同一输入永远同一答案。",
    fill=NAVY, color=WHITE, size=15, bold=True)
box(s, Inches(0.6), Inches(3.0), Inches(5.9), Inches(0.5), "对照题目", fill=GREEN, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(0.6), Inches(3.6), Inches(5.9), Inches(3.0), [
    "要求 ①：六级兜底阶梯 + 指纹冻结 + Agent —— 完成",
    "要求 ②：只用只读 CLI / syslog / trap，设备侧零改造 —— 完成",
    "要求 ③：拓扑上下文 + 跨设备取证 + 对照诊断 —— 完成",
    "判据 60% / 80% / 100%：真机实测 SSR 100%、字节一致 —— 完成",
], size=14)
box(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(0.5), "诚实的边界与后续", fill=ACCENT2, color=WHITE, bold=True, size=14, line=None)
bullets(s, Inches(6.9), Inches(3.6), Inches(5.9), Inches(3.0), [
    "F1 快照复用、F3 投票目前在单轮模式生效；Agent 模式可按「只对结论投票」接入",
    "命令预算在全网巡检类问题上贴近上限，可按规模调整",
    "MIB 里没有的 trap 按最长前缀诚实显示数字尾巴",
    "更多厂商只需导入手册 + MIB + 一条 vendors.json 预设",
], size=14)
num(s)

out = "docs/DetOps-答辩.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides))
